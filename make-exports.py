from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import io, os

slides_dir = "/Users/amar/Desktop/Limitra/slides-tmp"
out_dir    = "/Users/amar/Desktop/Limitra"
paths = [f"{slides_dir}/slide{i}.png" for i in range(6)]

# ── PDF ──────────────────────────────────────────────────────────────────────
images = [Image.open(p).convert("RGB") for p in paths]
images[0].save(
    f"{out_dir}/pitch-deck.pdf",
    save_all=True,
    append_images=images[1:],
)
print("PDF done")

# ── PPTX ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

for p in paths:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(p, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

prs.save(f"{out_dir}/pitch-deck.pptx")
print("PPTX done")
